import csv
import io

import docx
import fitz  # PyMuPDF
import openpyxl
import pytesseract
from pdf2image import convert_from_bytes, convert_from_path
from PIL import Image
from pptx import Presentation

from app.logger.logger import logger


class DocumentExtractor:
    def __init__(self):
        self.supported_extensions = {
            "pdf": self.extract_pdf,
            "doc": self.extract_word,
            "docx": self.extract_word,
            "xls": self.extract_excel,
            "xlsx": self.extract_excel,
            "ppt": self.extract_powerpoint,
            "pptx": self.extract_powerpoint,
            "csv": self.extract_csv,
            "txt": self.extract_text,
            "md": self.extract_text,
            "jpg": self.extract_image,
            "jpeg": self.extract_image,
            "png": self.extract_image,
            "bmp": self.extract_image,
            "gif": self.extract_image,
        }

    def extract(self, file_path):
        file_extension = file_path.split(".")[-1].lower()
        if file_extension not in self.supported_extensions:
            return "Unsupported file type"

        try:
            return self.supported_extensions[file_extension](file_path)
        except Exception as e:
            return f"Error processing file: {str(e)}"

    def extract_pdf(self, file_path):
        text = ""
        images = []
        doc = fitz.open(file_path)
        for page_num, page in enumerate(doc):
            # Extract text from the page
            page_text = page.get_text()
            if page_text.strip():
                text += page_text
            else:
                # If no text is found, render the page as an image and perform OCR
                try:
                    pix = page.get_pixmap()
                    image = Image.open(io.BytesIO(pix.tobytes()))
                    images.append(image)
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        text += f"\n[Page {page_num + 1}]:\n{ocr_text}"
                except Exception as e:
                    # Log the error and continue processing
                    logger.error(
                        f"Error during OCR of page image on page {page_num + 1}: {e}"
                    )
                    continue

            # Extract images from the page and perform OCR on them
            for img_index, img in enumerate(page.get_images()):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))
                    images.append(image)
                    ocr_text = pytesseract.image_to_string(image)
                    if ocr_text.strip():
                        text += f"\n[Page {page_num + 1}, Image {img_index + 1}]:\n{ocr_text}"
                except Exception as e:
                    # Log the error and continue processing
                    logger.error(
                        f"Error during OCR of image {img_index + 1} on page {page_num + 1}: {e}"
                    )
                    continue

        return {"text": text, "images": images}

    def extract_word(self, file_path):
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    def extract_excel(self, file_path):
        wb = openpyxl.load_workbook(file_path)
        data = {}
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            data[sheet_name] = [
                [cell.value for cell in row] for row in sheet.iter_rows()
            ]
        return data

    def extract_powerpoint(self, file_path):
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    def extract_csv(self, file_path):
        with open(file_path, "r") as file:
            reader = csv.reader(file)
            return list(reader)

    def extract_text(self, file_path):
        with open(file_path, "r") as file:
            return file.read()

    def extract_image(self, file_path):
        image = Image.open(file_path)
        text = pytesseract.image_to_string(image)
        return {"image": image, "text": text}
